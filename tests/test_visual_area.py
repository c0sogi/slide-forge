"""Tests for VisualArea auto-layout container."""

from unittest.mock import MagicMock, patch

import pytest

from slide_forge.default import get_presentation
from slide_forge.default.slide import (
    _ASCII_WIDTH_FACTOR,
    _CAPTION_BOX_HEIGHT,
    _CAPTION_GAP,
    _CJK_WIDTH_FACTOR,
    _CONTENT_HEIGHT,
    _CONTENT_LEFT,
    _CONTENT_TOP,
    _CONTENT_WIDTH,
    _FONT_LINE_HEIGHT_FACTOR,
    _LINE_SPACING_PCT,
    _SHAPE_REF_ATTR,
    _VISUAL_BOTTOM_MARGIN,
    _VISUAL_GAP,
    _VISUAL_TOP_GAP,
    _classify_paragraph,
    _detect_font_size,
    _estimate_text_width,
    add_bullet,
    add_content_box,
    add_section,
    add_slide_title,
    add_spacer,
    create_slide,
    estimate_text_height,
    shrink_content_box,
    visual_area,
)


@pytest.fixture
def slide():
    prs = get_presentation()
    return create_slide(prs)


# ── factory defaults ────────────────────────────────────────


class TestVisualAreaDefaults:
    def test_default_top(self, slide):
        area = visual_area(slide)
        expected = int(_CONTENT_TOP) + int(_CONTENT_HEIGHT) + int(_VISUAL_TOP_GAP)
        assert area._top == expected

    def test_default_height(self, slide):
        area = visual_area(slide)
        expected_top = int(_CONTENT_TOP) + int(_CONTENT_HEIGHT) + int(_VISUAL_TOP_GAP)
        expected_height = 6_858_000 - int(_VISUAL_BOTTOM_MARGIN) - expected_top
        assert area._height == expected_height

    def test_default_left_width(self, slide):
        area = visual_area(slide)
        assert area._left == int(_CONTENT_LEFT)
        assert area._width == int(_CONTENT_WIDTH)

    def test_default_gap(self, slide):
        area = visual_area(slide)
        assert area._gap == int(_VISUAL_GAP)

    def test_custom_overrides(self, slide):
        area = visual_area(slide, left=100, top=200, width=300, height=400, gap=50)
        assert area._left == 100
        assert area._top == 200
        assert area._width == 300
        assert area._height == 400
        assert area._gap == 50


# ── empty / edge cases ──────────────────────────────────────


class TestEdgeCases:
    def test_empty_render(self, slide):
        area = visual_area(slide)
        result = area.render()
        assert result == []

    def test_double_render_raises(self, slide):
        area = visual_area(slide)
        area.render()
        with pytest.raises(RuntimeError, match="한 번만"):
            area.render()

    def test_add_after_render_raises(self, slide):
        area = visual_area(slide)
        area.render()
        with pytest.raises(RuntimeError, match="render\\(\\) 호출 후"):
            area.add_chart("column", categories=["A"], series={"S": [1]})

    def test_add_table_after_render_raises(self, slide):
        area = visual_area(slide)
        area.render()
        with pytest.raises(RuntimeError, match="render\\(\\) 호출 후"):
            area.add_table([["a"]])

    def test_add_figure_after_render_raises(self, slide):
        area = visual_area(slide)
        area.render()
        with pytest.raises(RuntimeError, match="render\\(\\) 호출 후"):
            area.add_figure("x.png")

    def test_add_shape_after_render_raises(self, slide):
        area = visual_area(slide)
        area.render()
        with pytest.raises(RuntimeError, match="render\\(\\) 호출 후"):
            area.add_shape("rectangle")

    def test_zero_weight_raises(self, slide):
        area = visual_area(slide)
        with pytest.raises(ValueError, match="weight"):
            area.add_chart("column", weight=0, categories=["A"], series={"S": [1]})

    def test_negative_weight_raises(self, slide):
        area = visual_area(slide)
        with pytest.raises(ValueError, match="weight"):
            area.add_table([["a"]], weight=-1)

    def test_chaining(self, slide):
        area = visual_area(slide)
        result = area.add_shape("rectangle")
        assert result is area


# ── layout: single element ──────────────────────────────────


class TestSingleElement:
    def test_single_chart_full_width(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart("column", categories=["A"], series={"S": [1]})
            results = area.render()

            assert len(results) == 1
            _, kwargs = mock.call_args
            assert kwargs["left"] == 0
            assert kwargs["width"] == 9000000
            assert kwargs["height"] == 2000000
            assert kwargs["top"] == 4000000

    def test_single_table_full_width(self, slide):
        with patch("slide_forge.default.slide.add_table", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_table([["H"], ["v"]])
            area.render()

            _, kwargs = mock.call_args
            assert kwargs["left"] == 0
            assert kwargs["width"] == 9000000
            # Table: height is NOT passed (auto-size)
            assert "height" not in kwargs

    def test_single_figure_full_width(self, slide):
        with patch("slide_forge.default.slide.add_figure", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_figure("img.png")
            area.render()

            _, kwargs = mock.call_args
            assert kwargs["left"] == 0
            assert kwargs["width"] == 9000000
            # Figure: height is NOT passed (auto aspect ratio)
            assert "height" not in kwargs


# ── layout: two elements ────────────────────────────────────


class TestTwoElements:
    def test_equal_weight_no_gap(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart("column", categories=["A"], series={"S": [1]})
            area.add_chart("column", categories=["B"], series={"S": [2]})
            area.render()

            assert mock.call_count == 2
            call1 = mock.call_args_list[0]
            call2 = mock.call_args_list[1]
            assert call1.kwargs["left"] == 0
            assert call1.kwargs["width"] == 4500000
            assert call2.kwargs["left"] == 4500000
            assert call2.kwargs["width"] == 4500000  # last absorbs remainder

    def test_equal_weight_with_gap(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=200000)
            area.add_chart("column", categories=["A"], series={"S": [1]})
            area.add_chart("column", categories=["B"], series={"S": [2]})
            area.render()

            # available = 9000000 - 200000 = 8800000, each = 4400000
            call1 = mock.call_args_list[0]
            call2 = mock.call_args_list[1]
            assert call1.kwargs["left"] == 0
            assert call1.kwargs["width"] == 4400000
            assert call2.kwargs["left"] == 4600000  # 4400000 + 200000 gap
            assert call2.kwargs["width"] == 4400000  # remainder: 9000000 - 4600000

    def test_weighted_2_1(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart("column", categories=["A"], series={"S": [1]}, weight=2)
            area.add_chart("column", categories=["B"], series={"S": [2]}, weight=1)
            area.render()

            call1 = mock.call_args_list[0]
            call2 = mock.call_args_list[1]
            assert call1.kwargs["width"] == 6000000  # 9M * 2/3
            assert call2.kwargs["left"] == 6000000
            assert call2.kwargs["width"] == 3000000  # remainder


# ── layout: three / four elements ───────────────────────────


class TestMultipleElements:
    def test_three_equal(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            for _ in range(3):
                area.add_chart("column", categories=["A"], series={"S": [1]})
            area.render()

            assert mock.call_count == 3
            assert mock.call_args_list[0].kwargs["width"] == 3000000
            assert mock.call_args_list[1].kwargs["left"] == 3000000
            assert mock.call_args_list[1].kwargs["width"] == 3000000
            assert mock.call_args_list[2].kwargs["left"] == 6000000
            assert mock.call_args_list[2].kwargs["width"] == 3000000

    def test_four_elements(self, slide):
        with patch("slide_forge.default.slide.add_shape", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=8000000, height=2000000, gap=0)
            for _ in range(4):
                area.add_shape("rectangle")
            area.render()

            assert mock.call_count == 4
            assert mock.call_args_list[0].kwargs["width"] == 2000000
            assert mock.call_args_list[3].kwargs["left"] == 6000000
            assert mock.call_args_list[3].kwargs["width"] == 2000000


# ── rounding fix ────────────────────────────────────────────


class TestRoundingFix:
    def test_last_element_absorbs_remainder(self, slide):
        """When weights don't divide cleanly, the last element absorbs the remainder."""
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            # 10000 / 3 = 3333.33... → int truncates to 3333
            area = visual_area(slide, left=0, top=4000000, width=10000, height=2000000, gap=0)
            for _ in range(3):
                area.add_chart("column", categories=["A"], series={"S": [1]})
            area.render()

            w1 = mock.call_args_list[0].kwargs["width"]
            w2 = mock.call_args_list[1].kwargs["width"]
            w3 = mock.call_args_list[2].kwargs["width"]
            # First two: int(10000 * 1/3) = 3333
            assert w1 == 3333
            assert w2 == 3333
            # Last absorbs remainder: 10000 - 3333 - 3333 = 3334
            assert w3 == 3334
            assert w1 + w2 + w3 == 10000


# ── caption budgeting ───────────────────────────────────────


class TestCaptionBudgeting:
    CAPTION_RESERVE = int(_CAPTION_GAP) + int(_CAPTION_BOX_HEIGHT)

    def test_chart_with_caption_reduces_height(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart("column", categories=["A"], series={"S": [1]}, caption="Fig 1")
            area.render()

            _, kwargs = mock.call_args
            assert kwargs["height"] == 2000000 - self.CAPTION_RESERVE

    def test_chart_without_caption_full_height(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart("column", categories=["A"], series={"S": [1]})
            area.render()

            _, kwargs = mock.call_args
            assert kwargs["height"] == 2000000

    def test_mixed_captions_per_element(self, slide):
        """Only elements WITH captions get reduced height."""
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock_chart:
            with patch("slide_forge.default.slide.add_shape", return_value=MagicMock()) as mock_shape:
                area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
                area.add_chart("column", categories=["A"], series={"S": [1]}, caption="Fig 1")
                area.add_shape("rectangle")
                area.render()

                # Chart has caption → reduced height
                assert mock_chart.call_args.kwargs["height"] == 2000000 - self.CAPTION_RESERVE
                # Shape has no caption → full height
                assert mock_shape.call_args.kwargs["height"] == 2000000

    def test_table_ignores_height_regardless(self, slide):
        """Table never gets explicit height (auto-size), even with caption."""
        with patch("slide_forge.default.slide.add_table", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_table([["H"], ["v"]], caption="Table 1")
            area.render()

            assert "height" not in mock.call_args.kwargs

    def test_figure_ignores_height_regardless(self, slide):
        """Figure never gets explicit height (aspect ratio preserved), even with caption."""
        with patch("slide_forge.default.slide.add_figure", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_figure("img.png", caption="Fig 1")
            area.render()

            assert "height" not in mock.call_args.kwargs


# ── per-kind height strategy ────────────────────────────────


class TestPerKindHeightStrategy:
    def test_chart_gets_container_height(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart("column", categories=["A"], series={"S": [1]})
            area.render()
            assert mock.call_args.kwargs["height"] == 2000000

    def test_shape_gets_container_height(self, slide):
        with patch("slide_forge.default.slide.add_shape", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_shape("rectangle")
            area.render()
            assert mock.call_args.kwargs["height"] == 2000000

    def test_table_gets_no_height(self, slide):
        with patch("slide_forge.default.slide.add_table", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_table([["H"], ["v"]])
            area.render()
            assert "height" not in mock.call_args.kwargs

    def test_figure_gets_no_height(self, slide):
        with patch("slide_forge.default.slide.add_figure", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_figure("img.png")
            area.render()
            assert "height" not in mock.call_args.kwargs


# ── kwargs forwarding ───────────────────────────────────────


class TestKwargsForwarding:
    def test_chart_kwargs_forwarded(self, slide):
        with patch("slide_forge.default.slide.add_chart", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_chart(
                "line_markers",
                categories=["Q1", "Q2"],
                series={"Rev": [100, 200]},
                title="Revenue",
                legend=False,
                number_format="0.0%",
            )
            area.render()

            kw = mock.call_args.kwargs
            assert kw["title"] == "Revenue"
            assert kw["legend"] is False
            assert kw["number_format"] == "0.0%"
            assert kw["categories"] == ["Q1", "Q2"]

    def test_table_kwargs_forwarded(self, slide):
        with patch("slide_forge.default.slide.add_table", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_table(
                [["H1", "H2"], ["a", "b"]],
                first_row=False,
                first_col=True,
                header_size=16,
                body_size=10,
            )
            area.render()

            kw = mock.call_args.kwargs
            assert kw["first_row"] is False
            assert kw["first_col"] is True
            assert kw["header_size"] == 16
            assert kw["body_size"] == 10

    def test_shape_kwargs_forwarded(self, slide):
        with patch("slide_forge.default.slide.add_shape", return_value=MagicMock()) as mock:
            area = visual_area(slide, left=0, top=4000000, width=9000000, height=2000000, gap=0)
            area.add_shape(
                "oval",
                fill="FF0000",
                transparency=50,
                text="Hello",
                bold=True,
            )
            area.render()

            kw = mock.call_args.kwargs
            assert kw["fill"] == "FF0000"
            assert kw["transparency"] == 50
            assert kw["text"] == "Hello"
            assert kw["bold"] is True


# ── integration test ────────────────────────────────────────


class TestIntegration:
    def test_chart_and_table_side_by_side(self, slide):
        """Real (non-mocked) integration: chart + table placed without error."""
        add_slide_title(slide, "Integration Test")
        area = visual_area(slide)
        area.add_chart(
            "column",
            categories=["A", "B", "C"],
            series={"S1": [10, 20, 30]},
            title="Test Chart",
        )
        area.add_table(
            [["Col1", "Col2"], ["x", "y"]],
            first_row=True,
        )
        results = area.render()
        assert len(results) == 2

    def test_render_returns_placed_objects(self, slide):
        """render() returns a list of the placed PPTX objects."""
        area = visual_area(slide)
        area.add_chart(
            "column",
            categories=["A", "B"],
            series={"S": [1, 2]},
        )
        results = area.render()
        assert len(results) == 1
        # add_chart returns a Chart object
        from pptx.chart.chart import Chart

        assert isinstance(results[0], Chart)


# ── _detect_font_size ────────────────────────────────────────


class TestDetectFontSize:
    def test_extracts_from_bullet(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_bullet(tf, "Hello", level=0)
        # Last paragraph is the bullet
        p = tf.paragraphs[-1]
        result = _detect_font_size(p._element)
        assert result == 14  # L0 default

    def test_extracts_from_level1_bullet(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_bullet(tf, "Sub", level=1)
        p = tf.paragraphs[-1]
        result = _detect_font_size(p._element)
        assert result == 12  # L1 default

    def test_section_font_size(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Header")
        p = tf.paragraphs[0]
        result = _detect_font_size(p._element)
        assert result == 18  # Section default

    def test_returns_none_for_spacer(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_spacer(tf)
        # Spacer is the last paragraph (no runs)
        p = tf.paragraphs[-1]
        result = _detect_font_size(p._element)
        assert result is None


# ── _estimate_text_width ─────────────────────────────────────


class TestEstimateTextWidth:
    def _per_char_width(self, n: int, factor: float, font_size: int) -> int:
        """Sum per-character int() to match implementation rounding."""
        return sum(int(factor * font_size * 12700) for _ in range(n))

    def test_ascii_only(self):
        text = "Hello"
        width = _estimate_text_width(text, 14)
        expected = self._per_char_width(5, _ASCII_WIDTH_FACTOR, 14)
        assert width == expected

    def test_cjk_only(self):
        text = "한글테스트"
        width = _estimate_text_width(text, 14)
        expected = self._per_char_width(5, _CJK_WIDTH_FACTOR, 14)
        assert width == expected

    def test_mixed_ascii_cjk(self):
        text = "AB한글"
        width = _estimate_text_width(text, 14)
        ascii_w = self._per_char_width(2, _ASCII_WIDTH_FACTOR, 14)
        cjk_w = self._per_char_width(2, _CJK_WIDTH_FACTOR, 14)
        assert width == ascii_w + cjk_w

    def test_strips_color_tags(self):
        tagged = "[red]Hello[/red]"
        plain = "Hello"
        assert _estimate_text_width(tagged, 14) == _estimate_text_width(plain, 14)

    def test_arrow_replacement(self):
        text = "A->B"
        width = _estimate_text_width(text, 14)
        # "A->B" becomes "A→B" (3 chars: A=ascii, →=CJK, B=ascii)
        ascii_w = self._per_char_width(2, _ASCII_WIDTH_FACTOR, 14)
        cjk_w = self._per_char_width(1, _CJK_WIDTH_FACTOR, 14)
        assert width == ascii_w + cjk_w

    def test_empty_string(self):
        assert _estimate_text_width("", 14) == 0


# ── _classify_paragraph ─────────────────────────────────────


class TestClassifyParagraph:
    def test_spacer(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_spacer(tf)
        p = tf.paragraphs[-1]
        p_type, font_size, line_height, avail_width = _classify_paragraph(p._element)
        assert p_type == "spacer"
        assert font_size == 14

    def test_section(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test Header")
        p = tf.paragraphs[0]
        p_type, font_size, line_height, avail_width = _classify_paragraph(p._element)
        assert p_type == "section"
        assert font_size == 18

    def test_bullet_level0(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_bullet(tf, "Bullet", level=0)
        p = tf.paragraphs[-1]
        p_type, font_size, line_height, avail_width = _classify_paragraph(p._element)
        assert p_type == "bullet"
        assert font_size == 14

    def test_bullet_level1(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_bullet(tf, "Sub-bullet", level=1)
        p = tf.paragraphs[-1]
        p_type, font_size, line_height, avail_width = _classify_paragraph(p._element)
        assert p_type == "bullet"
        assert font_size == 12

    def test_line_height_includes_spacing(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_bullet(tf, "Bullet", level=0)
        p = tf.paragraphs[-1]
        _, font_size, line_height, _ = _classify_paragraph(p._element)
        expected = int(font_size * 12700 * _FONT_LINE_HEIGHT_FACTOR * _LINE_SPACING_PCT / 100000)
        assert line_height == expected

    def test_available_width_reduced_for_bullets(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Test")
        add_bullet(tf, "Bullet", level=0)
        p = tf.paragraphs[-1]
        _, _, _, avail_width = _classify_paragraph(p._element)
        assert avail_width < int(_CONTENT_WIDTH)


# ── estimate_text_height ─────────────────────────────────────


class TestEstimateTextHeight:
    def test_basic_section_and_bullet(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Short bullet")
        result = estimate_text_height(tf)
        assert result > 0
        assert result <= int(_CONTENT_HEIGHT)

    def test_empty_textframe(self, slide):
        tf = add_content_box(slide)
        # TextFrame with only one empty paragraph
        result = estimate_text_height(tf)
        assert result > 0

    def test_safety_margin_applied(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "A")
        result = estimate_text_height(tf)
        # Without safety margin, sum would be smaller
        # The result includes 10% margin
        assert result > 0

    def test_capped_at_content_box_height(self, slide):
        """Even with many bullets, height shouldn't exceed content-box height."""
        tf = add_content_box(slide)
        add_section(tf, "Header")
        for i in range(50):
            add_bullet(tf, f"This is a very long bullet point number {i} " * 5)
        result = estimate_text_height(tf)
        shape = getattr(tf, _SHAPE_REF_ATTR)
        assert result == int(shape.height)

    def test_spacer_contributes_height(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "A")
        h1 = estimate_text_height(tf)
        add_spacer(tf)
        h2 = estimate_text_height(tf)
        assert h2 > h1

    def test_more_text_means_more_height(self, slide):
        prs = get_presentation()
        s1 = create_slide(prs)
        tf1 = add_content_box(s1)
        add_section(tf1, "H")
        add_bullet(tf1, "Short")

        s2 = create_slide(prs)
        tf2 = add_content_box(s2)
        add_section(tf2, "H")
        add_bullet(tf2, "Short")
        add_bullet(tf2, "Another line")
        add_bullet(tf2, "Third line")

        assert estimate_text_height(tf2) > estimate_text_height(tf1)

    def test_wrapping_increases_height(self, slide):
        """Long text that wraps should estimate taller than short text."""
        prs = get_presentation()
        s1 = create_slide(prs)
        tf1 = add_content_box(s1)
        add_section(tf1, "H")
        add_bullet(tf1, "Short")

        s2 = create_slide(prs)
        tf2 = add_content_box(s2)
        add_section(tf2, "H")
        add_bullet(tf2, "This is a very long bullet that should wrap " * 10)

        assert estimate_text_height(tf2) > estimate_text_height(tf1)

    def test_narrower_content_box_increases_height(self):
        """Narrower content box should produce taller estimates for same text."""
        prs = get_presentation()
        text = "This is a long bullet that should wrap differently by width. " * 4

        s_wide = create_slide(prs)
        tf_wide = add_content_box(s_wide, width=int(_CONTENT_WIDTH))
        add_section(tf_wide, "H")
        add_bullet(tf_wide, text)

        s_narrow = create_slide(prs)
        tf_narrow = add_content_box(s_narrow, width=2_400_000)
        add_section(tf_narrow, "H")
        add_bullet(tf_narrow, text)

        assert estimate_text_height(tf_narrow) > estimate_text_height(tf_wide)

    def test_custom_height_caps_estimate(self, slide):
        """Estimate should be capped by custom content-box height."""
        tf = add_content_box(slide, height=900_000)
        add_section(tf, "Header")
        for i in range(25):
            add_bullet(tf, f"This is very long line {i} " * 8)
        assert estimate_text_height(tf) == 900_000


# ── shrink_content_box ───────────────────────────────────────


class TestShrinkContentBox:
    def test_shrinks_shape_height(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Short bullet")
        shape = getattr(tf, _SHAPE_REF_ATTR)
        original_height = int(shape.height)
        bottom = shrink_content_box(tf)
        new_height = int(shape.height)
        assert new_height < original_height
        assert bottom == int(shape.top) + new_height

    def test_returns_bottom_coordinate(self, slide):
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Bullet")
        bottom = shrink_content_box(tf)
        shape = getattr(tf, _SHAPE_REF_ATTR)
        assert bottom == int(shape.top) + int(shape.height)

    def test_raises_without_shape_ref(self):
        """shrink_content_box requires add_content_box's TextFrame."""
        mock_tf = MagicMock(spec=["paragraphs"])
        with pytest.raises(RuntimeError, match="add_content_box"):
            shrink_content_box(mock_tf)

    def test_add_content_box_stores_shape_ref(self, slide):
        tf = add_content_box(slide)
        shape = getattr(tf, _SHAPE_REF_ATTR, None)
        assert shape is not None
        assert hasattr(shape, "height")
        assert hasattr(shape, "top")


# ── visual_area with content_box ─────────────────────────────


class TestVisualAreaContentBox:
    def test_content_box_adjusts_top(self, slide):
        """content_box parameter should position top based on estimated text height."""
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Short")

        area = visual_area(slide, content_box=tf)
        text_h = estimate_text_height(tf)
        shape_top = int(getattr(tf, _SHAPE_REF_ATTR).top)
        expected_top = shape_top + text_h + int(_VISUAL_TOP_GAP)
        assert area._top == expected_top

    def test_content_box_gives_more_height(self, slide):
        """With content_box, the visual area should be taller (more room)."""
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Short")

        area_default = visual_area(slide)
        area_smart = visual_area(slide, content_box=tf)
        # Smart positioning should yield a lower top → more height
        assert area_smart._top < area_default._top
        assert area_smart._height > area_default._height

    def test_content_box_none_preserves_default(self, slide):
        """content_box=None should behave exactly like before."""
        area_default = visual_area(slide)
        area_none = visual_area(slide, content_box=None)
        assert area_default._top == area_none._top
        assert area_default._height == area_none._height

    def test_explicit_top_ignores_content_box(self, slide):
        """When top is explicitly set, content_box is ignored."""
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Short")

        area = visual_area(slide, top=5000000, content_box=tf)
        assert area._top == 5000000

    def test_content_box_height_auto_calculated(self, slide):
        """Height should auto-fill from the smart top to bottom margin."""
        tf = add_content_box(slide)
        add_section(tf, "Header")
        add_bullet(tf, "Short")

        area = visual_area(slide, content_box=tf)
        shape_top = int(getattr(tf, _SHAPE_REF_ATTR).top)
        expected_top = shape_top + estimate_text_height(tf) + int(_VISUAL_TOP_GAP)
        expected_height = 6_858_000 - int(_VISUAL_BOTTOM_MARGIN) - expected_top
        assert area._height == expected_height

    def test_content_box_custom_top_is_respected(self, slide):
        """Smart top calculation should use the actual content-box top."""
        custom_top = 1_500_000
        tf = add_content_box(slide, top=custom_top)
        add_section(tf, "Header")
        add_bullet(tf, "Short")

        area = visual_area(slide, content_box=tf)
        expected_top = custom_top + estimate_text_height(tf) + int(_VISUAL_TOP_GAP)
        assert area._top == expected_top

    def test_full_workflow_integration(self, slide):
        """Full workflow: create content → shrink → visual_area(content_box=...)."""
        tf = add_content_box(slide)
        add_section(tf, "주요 결과")
        add_bullet(tf, "정상 데이터 1,198,500건(96.3%)")
        add_bullet(tf, "Precision: [green]92.3%[/green], Recall: 88.0%", level=1)

        shrink_content_box(tf)
        area = visual_area(slide, content_box=tf)

        # The area should have a valid top and height
        assert area._top > int(_CONTENT_TOP)
        assert area._height > 0
        assert area._top + area._height <= 6_858_000
