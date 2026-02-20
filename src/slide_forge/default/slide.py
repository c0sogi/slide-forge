import re
from typing import Sequence

from lxml import etree
from pptx.chart.chart import Chart
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.presentation import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pptx.table import Table
from pptx.text.text import TextFrame, _Paragraph
from pptx.util import Emu, Pt

# ── Default fonts ─────────────────────────────────────────────
FONT_ASCII = "Arial"
FONT_EA = "HY\uacac\uace0\ub515"  # HY견고딕

# ── Layout / Master indices ───────────────────────────────────
_BLANK_LAYOUT_NAME = "\ube48 \ud654\uba74"  # 빈 화면
_COVER_MASTER_IDX = 0  # Master 0: full background image (앞/뒤 표지)
_CONTENT_MASTER_IDX = 1  # Master 1: header lines + logos (내용)

# ── Slide type tag ───────────────────────────────────────────
_SLIDE_TYPE_ATTR = "_slide_forge_type"
_COVER = "cover"
_CONTENT = "content"

# ── Section state tracking ───────────────────────────────────
_SECTION_ADDED_ATTR = "_slide_forge_section_added"

# ── Cover slide ──────────────────────────────────────────────
_COVER_TITLE_LEFT = Emu(561257)
_COVER_TITLE_TOP = Emu(2636912)
_COVER_TITLE_WIDTH = Emu(8010140)
_COVER_TITLE_HEIGHT = Emu(1186800)
_COVER_TITLE_SIZE = 20
_COVER_COLOR = RGBColor(0x00, 0x20, 0x60)

_COVER_INFO_LEFT = Emu(3242805)
_COVER_INFO_TOP = Emu(5229200)
_COVER_INFO_WIDTH = Emu(5328592)
_COVER_INFO_HEIGHT = Emu(696857)
_COVER_INFO_SIZE = 14

# ── Slide title ───────────────────────────────────────────────
_TITLE_LEFT = Emu(179512)
_TITLE_TOP = Emu(154732)
_TITLE_WIDTH = Emu(8352928)
_TITLE_HEIGHT = Emu(461665)
_TITLE_SIZE = 24
_TITLE_COLOR = RGBColor(0x07, 0x2A, 0x5E)

# ── Content text box ─────────────────────────────────────────
_CONTENT_LEFT = Emu(166261)
_CONTENT_TOP = Emu(724090)
_CONTENT_WIDTH = Emu(8870234)
_CONTENT_HEIGHT = Emu(3599447)

# ── Line spacing (130 %) ─────────────────────────────────────
_LINE_SPACING_PCT = 130000

# ── Section header (▌) ──────────────────────────────────────
_SECTION_SIZE = 18
_SECTION_COLOR = RGBColor(0x40, 0x40, 0x40)

# ── Bullet margin ────────────────────────────────────────────
_MARGIN_BASE = 266700
_MARGIN_STEP = 182563  # per-level increment

# ── Per-level defaults ───────────────────────────────────────
_L0_INDENT = "-174625"
_L0_SIZE = 14
_L0_COLOR = RGBColor(0x40, 0x40, 0x40)

_L1_INDENT = "-182563"
_L1_SIZE = 12
_L1_COLOR = RGBColor(0x40, 0x40, 0x40)

# ── Bullet styles ────────────────────────────────────────────
_BULLET_STYLES: dict[str, dict[str, str]] = {
    "dash": {"char": "-", "font": FONT_EA, "pitch": "18", "charset": "-127"},
    "arrow": {"char": "\u00d8", "font": "Wingdings", "pitch": "2", "charset": "2"},
    "square": {"char": "n", "font": "Wingdings", "pitch": "2", "charset": "2"},
    "circle": {"char": "l", "font": "Wingdings", "pitch": "2", "charset": "2"},
    "check": {"char": "\u00fc", "font": "Wingdings", "pitch": "2", "charset": "2"},
}

_DEFAULT_BULLET = {0: "dash", 1: "arrow", 2: "square", 3: "circle", 4: "check"}

# ── Caption ──────────────────────────────────────────────────
_CAPTION_SIZE = 10
_CAPTION_COLOR = RGBColor(0x7F, 0x7F, 0x7F)
_CAPTION_GAP = Emu(50000)  # ~4 pt gap above caption
_CAPTION_BOX_HEIGHT = Emu(230000)  # caption text box height

# ── Inline color tags ────────────────────────────────────────
_COLOR_TAG_RE = re.compile(r"\[(#[0-9A-Fa-f]{6}|\w+)\](.*?)\[/\1\]")

_COLOR_MAP: dict[str, RGBColor] = {
    "red": RGBColor(0xC0, 0x00, 0x00),
    "green": RGBColor(0x00, 0xB0, 0x50),
    "blue": RGBColor(0x00, 0x70, 0xC0),
    "orange": RGBColor(0xFF, 0x7F, 0x00),
    "purple": RGBColor(0x70, 0x30, 0xA0),
    "navy": RGBColor(0x07, 0x2A, 0x5E),
    "teal": RGBColor(0x00, 0x80, 0x80),
    "gray": RGBColor(0x80, 0x80, 0x80),
    "grey": RGBColor(0x80, 0x80, 0x80),
}

# ── Color type ──────────────────────────────────────────────
Color = RGBColor | str
"""Accepted colour formats: RGBColor, ``"#RRGGBB"``, ``"RRGGBB"``, or named (``"red"``, ``"navy"``, etc.)."""


def _to_rgb(value: "Color") -> RGBColor:
    """Convert a *Color* value to :class:`RGBColor`.

    Accepts :class:`RGBColor` (pass-through), ``"#RRGGBB"`` or ``"RRGGBB"``
    hex strings, ``"R,G,B"`` decimal strings, and named colours from
    ``_COLOR_MAP``.
    """
    if isinstance(value, RGBColor):
        return value
    if not isinstance(value, str):
        raise TypeError(f"Expected RGBColor or str, got {type(value).__name__}")
    named = _COLOR_MAP.get(value.lower())
    if named is not None:
        return named
    if "," in value:
        parts = [p.strip() for p in value.split(",")]
        if len(parts) == 3:
            try:
                r, g, b = (int(p) for p in parts)
                return RGBColor(r, g, b)
            except (ValueError, OverflowError):
                pass
    h = value.lstrip("#")
    if len(h) == 6:
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            pass
    raise ValueError(f"Cannot parse colour: {value!r}")


# ── Arrow token ──────────────────────────────────────────────
_ARROW_CHAR = "\uf0e0"  # Wingdings right arrow (→)

# ── Shape types ─────────────────────────────────────────────
_SHAPE_TYPE_MAP: dict[str, MSO_SHAPE] = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "hexagon": MSO_SHAPE.HEXAGON,
    "chevron": MSO_SHAPE.CHEVRON,
}

# ── Line dash styles ───────────────────────────────────────
_DASH_MAP: dict[str, MSO_LINE_DASH_STYLE] = {
    "solid": MSO_LINE_DASH_STYLE.SOLID,
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
    "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
    "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
}

# ── Vertical alignment for shapes ──────────────────────────
_VALIGN_XML: dict[str, str] = {
    "top": "t",
    "middle": "ctr",
    "bottom": "b",
}


# ─── internal helpers ────────────────────────────────────────


def _resolve_color(name: str) -> RGBColor | None:
    """Resolve a color name or ``#RRGGBB`` hex code to :class:`RGBColor`."""
    if name.startswith("#") and len(name) == 7:
        try:
            return RGBColor(int(name[1:3], 16), int(name[3:5], 16), int(name[5:7], 16))
        except ValueError:
            return None
    return _COLOR_MAP.get(name.lower())


def _parse_color_tags(text: str) -> list[tuple[str, RGBColor | None]]:
    """Parse ``[color]text[/color]`` inline markup.

    Returns ``[(text, color_override), ...]``.
    *color_override* is ``None`` for untagged portions.
    Unknown color names are left as literal text (tags preserved).
    """
    segments: list[tuple[str, RGBColor | None]] = []
    last_end = 0
    for m in _COLOR_TAG_RE.finditer(text):
        color = _resolve_color(m.group(1))
        if color is None:
            continue  # unknown color — keep tags as literal text
        if m.start() > last_end:
            segments.append((text[last_end : m.start()], None))
        segments.append((m.group(2), color))
        last_end = m.end()
    if last_end < len(text):
        segments.append((text[last_end:], None))
    if not segments:
        segments.append((text, None))
    return segments


def _is_ascii(ch: str) -> bool:
    return bool(ch) and ord(ch) < 128


def _split_by_script(text: str) -> list[tuple[str, bool]]:
    """Split *text* into ``(segment, is_ascii)`` chunks.

    Whitespace is kept with the preceding chunk so that run boundaries
    don't land on spaces.
    """
    if not text:
        return []

    chunks: list[tuple[str, bool]] = []
    buf: list[str] = [text[0]]
    cur = _is_ascii(text[0])

    for ch in text[1:]:
        if ch.isspace():
            buf.append(ch)
            continue
        a = _is_ascii(ch)
        if a == cur:
            buf.append(ch)
        else:
            chunks.append(("".join(buf), cur))
            buf = [ch]
            cur = a

    chunks.append(("".join(buf), cur))
    return chunks


def _set_ea_font(run, typeface: str = FONT_EA) -> None:
    """Append ``<a:ea typeface="…"/>`` to a run's ``rPr``."""
    rPr = run._r.get_or_add_rPr()
    for existing in rPr.findall(qn("a:ea")):
        rPr.remove(existing)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", typeface)


def _set_sym_font(run, typeface: str = "Wingdings") -> None:
    """Append ``<a:sym typeface="…"/>`` so ``\\uf0e0`` renders as →."""
    rPr = run._r.get_or_add_rPr()
    for existing in rPr.findall(qn("a:sym")):
        rPr.remove(existing)
    sym = etree.SubElement(rPr, qn("a:sym"))
    sym.set("typeface", typeface)
    sym.set("pitchFamily", "2")
    sym.set("charset", "2")


def _set_line_spacing(pPr, pct: int = _LINE_SPACING_PCT) -> None:
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", str(pct))


def _set_bullet(
    pPr,
    char: str,
    typeface: str,
    pitch: str,
    charset: str,
) -> None:
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", typeface)
    buFont.set("pitchFamily", pitch)
    buFont.set("charset", charset)
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)


def _set_auto_number(pPr, num_type: str = "arabicPeriod") -> None:
    """Add ``<a:buAutoNum type="…"/>`` for numbered bullets."""
    buAutoNum = etree.SubElement(pPr, qn("a:buAutoNum"))
    buAutoNum.set("type", num_type)


def _add_runs(
    p: _Paragraph,
    text: str,
    *,
    size: int,
    color: "Color",
    ascii_font: str = FONT_ASCII,
    ea_font: str = FONT_EA,
    ascii_bold: bool = True,
    unicode_bold: bool = False,
) -> None:
    """Add mixed-font runs to *p*, splitting on ASCII / non-ASCII.

    Supports ``[color]text[/color]`` inline markup for per-span color
    overrides.  Named colors (``red``, ``green``, …) and hex codes
    (``#RRGGBB``) are both accepted.

    ``->`` is automatically replaced with the Wingdings right-arrow
    token (``\\uf0e0``).
    """
    color = _to_rgb(color)
    text = text.replace("->", _ARROW_CHAR)
    for seg_text, color_override in _parse_color_tags(text):
        effective_color = color_override if color_override is not None else color
        for chunk, is_ascii in _split_by_script(seg_text):
            run = p.add_run()
            run.text = chunk
            run.font.name = ascii_font if is_ascii else ea_font
            run.font.bold = ascii_bold if is_ascii else unicode_bold
            run.font.size = Pt(size)
            run.font.color.rgb = effective_color
            _set_ea_font(run, ea_font)
            if _ARROW_CHAR in chunk:
                _set_sym_font(run)


def _next_para(tf: TextFrame) -> _Paragraph:
    """Re-use the empty first paragraph or append a new one."""
    paras = tf.paragraphs
    if len(paras) == 1 and not paras[0].text and not paras[0].runs:
        return paras[0]
    return tf.add_paragraph()


# ─── internal: slide type validation ─────────────────────────


def _require(slide: Slide, expected: str) -> None:
    """Raise :class:`TypeError` if *slide* is not the expected type."""
    actual = getattr(slide, _SLIDE_TYPE_ATTR, None)
    if actual is None:
        raise TypeError("이 슬라이드는 create_slide() 또는 create_cover_slide()로 생성되지 않았습니다.")
    if actual == expected:
        return
    if expected == _COVER:
        raise TypeError(
            "이 함수는 표지 슬라이드(create_cover_slide)에서만 사용할 수 있습니다.\n"
            "  내용 슬라이드 → add_slide_title(), add_content_box() 등을 사용하세요."
        )
    raise TypeError(
        "이 함수는 내용 슬라이드(create_slide)에서만 사용할 수 있습니다.\n"
        "  표지 슬라이드 → add_cover_title(), add_cover_info()를 사용하세요."
    )


# ─── internal: layout lookup ─────────────────────────────────


def _get_blank_layout(prs: Presentation, master_idx: int):
    """Return the '빈 화면' layout from a specific slide master."""
    master = prs.slide_masters[master_idx]
    for layout in master.slide_layouts:
        if layout.name == _BLANK_LAYOUT_NAME:
            return layout
    raise ValueError(f"Layout '{_BLANK_LAYOUT_NAME}' not found in master {master_idx}")


# ─── public API ──────────────────────────────────────────────


def create_slide(prs: Presentation) -> Slide:
    """Add a **content** slide (Master 1: header lines + logos)."""
    slide = prs.slides.add_slide(_get_blank_layout(prs, _CONTENT_MASTER_IDX))
    setattr(slide, _SLIDE_TYPE_ATTR, _CONTENT)
    return slide


def create_cover_slide(prs: Presentation) -> Slide:
    """Add a **cover** slide (Master 0: full background image, 앞/뒤 표지)."""
    slide = prs.slides.add_slide(_get_blank_layout(prs, _COVER_MASTER_IDX))
    setattr(slide, _SLIDE_TYPE_ATTR, _COVER)
    return slide


def add_slide_title(
    slide: Slide,
    text: str,
    *,
    left: int = _TITLE_LEFT,
    top: int = _TITLE_TOP,
    width: int = _TITLE_WIDTH,
    height: int = _TITLE_HEIGHT,
    font_size: int = _TITLE_SIZE,
    color: "Color" = _TITLE_COLOR,
) -> None:
    """Add the top-of-slide title text box (content slides only)."""
    _require(slide, _CONTENT)
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _add_runs(p, text, size=font_size, color=color)


def add_content_box(
    slide: Slide,
    *,
    left: int = _CONTENT_LEFT,
    top: int = _CONTENT_TOP,
    width: int = _CONTENT_WIDTH,
    height: int = _CONTENT_HEIGHT,
) -> TextFrame:
    """Create the main content area and return its :class:`TextFrame`
    (content slides only).

    Pass the returned object to :func:`add_section`, :func:`add_bullet`,
    and :func:`add_spacer` to populate the slide.
    """
    _require(slide, _CONTENT)
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.clear()
    return tf


def _require_section(tf: TextFrame) -> None:
    """Raise :class:`RuntimeError` if ``add_section()`` has not been called on *tf* yet."""
    if not getattr(tf, _SECTION_ADDED_ATTR, False):
        raise RuntimeError(
            "add_section()을 먼저 호출해야 합니다.\n"
            "  올바른 순서: add_content_box() → add_section() → add_bullet() / add_spacer()"
        )


def add_section(
    tf: TextFrame,
    title: str,
    *,
    font_size: int = _SECTION_SIZE,
    color: "Color" = _SECTION_COLOR,
) -> _Paragraph:
    """Add a **▌Section Header** paragraph.

    Must be called exactly once per :class:`TextFrame` returned by
    :func:`add_content_box`, and before any :func:`add_bullet` or
    :func:`add_spacer` calls.
    """
    if getattr(tf, _SECTION_ADDED_ATTR, False):
        raise RuntimeError(
            "add_section()은 TextFrame당 한 번만 호출할 수 있습니다.\n"
            "  새 섹션이 필요하면 새 슬라이드를 만드세요: create_slide() → add_content_box() → add_section()"
        )
    setattr(tf, _SECTION_ADDED_ATTR, True)
    p = _next_para(tf)
    pPr = p._element.get_or_add_pPr()
    pPr.set("eaLnBrk", "1")
    pPr.set("hangingPunct", "1")
    _set_line_spacing(pPr, _LINE_SPACING_PCT)
    _add_runs(p, "\u258c" + title, size=font_size, color=color)
    return p


def add_bullet(
    tf: TextFrame,
    text: str,
    *,
    level: int = 0,
    bullet: str | None = None,
    font_size: int | None = None,
    color: "Color | None" = None,
) -> _Paragraph:
    """Add a bullet paragraph.

    Parameters
    ----------
    level : int
        Indentation depth.  ``0`` is top-level, ``1+`` is nested.
        Margin increases by ``_MARGIN_STEP`` per level.
    bullet : str, optional
        Bullet style name.  Built-in styles:

        * ``"dash"``   — ``-`` (HY견고딕)  *default for level 0*
        * ``"arrow"``  — ``Ø`` → (Wingdings)  *default for level 1*
        * ``"square"`` — ``■`` (Wingdings)  *default for level 2*
        * ``"circle"`` — ``●`` (Wingdings)  *default for level 3*
        * ``"check"``  — ``✓`` (Wingdings)  *default for level 4+*
        * ``"number"`` — auto-numbered (1. 2. 3.)
        * ``"none"``   — no bullet marker

        If *None*, defaults to ``"dash"`` for level 0 and ``"arrow"``
        otherwise.
    color : RGBColor, optional
        Override the default ``#404040``.
    """
    _require_section(tf)

    if bullet is None:
        bullet = _DEFAULT_BULLET.get(level, "check")

    indent = _L0_INDENT if level == 0 else _L1_INDENT
    sz = font_size or (_L0_SIZE if level == 0 else _L1_SIZE)
    clr = color or (_L0_COLOR if level == 0 else _L1_COLOR)
    marl = str(_MARGIN_BASE + level * _MARGIN_STEP)

    p = _next_para(tf)
    pPr = p._element.get_or_add_pPr()
    pPr.set("marL", marl)
    pPr.set("indent", indent)
    pPr.set("eaLnBrk", "1")
    pPr.set("hangingPunct", "1")
    _set_line_spacing(pPr, _LINE_SPACING_PCT)

    if bullet == "number":
        _set_auto_number(pPr)
    elif bullet != "none" and bullet in _BULLET_STYLES:
        s = _BULLET_STYLES[bullet]
        _set_bullet(pPr, s["char"], s["font"], s["pitch"], s["charset"])

    _add_runs(p, text, size=sz, color=clr)
    return p


def add_spacer(tf: TextFrame) -> _Paragraph:
    """Add an empty separator line."""
    _require_section(tf)
    p = _next_para(tf)
    pPr = p._element.get_or_add_pPr()
    pPr.set("marL", str(_MARGIN_BASE))
    _set_line_spacing(pPr, _LINE_SPACING_PCT)
    return p


def add_caption(
    slide: Slide,
    text: str,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int = _CAPTION_SIZE,
    color: "Color | None" = None,
) -> None:
    """Add a small caption text box (typically under a figure)."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    clr = color or _CAPTION_COLOR
    _add_runs(p, text, size=font_size, color=clr)


def _add_caption_below(
    slide: Slide,
    text: str,
    *,
    left: int,
    bottom: int,
    width: int,
) -> None:
    """Add a caption text box just below an element.

    *bottom* is the Y-coordinate (EMU) of the element's bottom edge.
    """
    txBox = slide.shapes.add_textbox(Emu(left), Emu(bottom + int(_CAPTION_GAP)), Emu(width), _CAPTION_BOX_HEIGHT)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_runs(p, text, size=_CAPTION_SIZE, color=_CAPTION_COLOR)


# ─── shape helpers ──────────────────────────────────────────


def _apply_fill_transparency(spPr, transparency: int) -> None:
    """Set fill transparency (0–100) via alpha on the ``srgbClr`` element."""
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is None:
        return
    alpha_el = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha_el.set("val", str((100 - transparency) * 1000))


def _apply_shadow(spPr, shadow: dict) -> None:
    """Apply an outer shadow effect to shape properties.

    *shadow* keys:

    * ``blur``    – blur radius in **points** (default 3)
    * ``offset``  – distance in **points** (default 2)
    * ``angle``   – direction in degrees 0–359 (default 135 = bottom-right)
    * ``color``   – :class:`RGBColor` (default black)
    * ``opacity`` – 0.0–1.0 (default 0.15)
    """
    blur = int(shadow.get("blur", 3) * 12700)  # pt → EMU
    offset = int(shadow.get("offset", 2) * 12700)  # pt → EMU
    angle = shadow.get("angle", 135)
    color = shadow.get("color", RGBColor(0, 0, 0))
    if isinstance(color, str):
        color = _to_rgb(color)
    opacity = shadow.get("opacity", 0.15)

    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)

    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", str(blur))
    outerShdw.set("dist", str(offset))
    outerShdw.set("dir", str(int(angle * 60000)))
    outerShdw.set("algn", "ctr")
    outerShdw.set("rotWithShape", "0")

    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgbClr.set("val", str(color))
    alpha_el = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha_el.set("val", str(int(opacity * 100000)))


# ─── shape / line (public) ──────────────────────────────────


def add_shape(
    slide: Slide,
    shape_type: str,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    fill: "Color | None" = None,
    transparency: int = 0,
    line_color: "Color | None" = None,
    line_width: int = 1,
    line_dash: str | None = None,
    shadow: dict | None = None,
    text: str | None = None,
    font_size: int = 12,
    color: "Color | None" = None,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    valign: str = "middle",
):
    """Add a shape to *slide*.

    Parameters
    ----------
    shape_type : str
        ``"rectangle"``, ``"oval"``, ``"rounded_rectangle"``,
        ``"triangle"``, ``"diamond"``, ``"hexagon"``, ``"chevron"``.
    left, top, width, height : int
        Position and size in EMU.
    fill : RGBColor, optional
        Solid fill colour.  ``None`` → no fill (transparent).
    transparency : int
        Fill transparency 0–100 (default 0 = fully opaque).
    line_color : RGBColor, optional
        Border colour.  ``None`` → no border.
    line_width : int
        Border width in EMU (default 1 pt).
    line_dash : str, optional
        ``"solid"``, ``"dash"``, ``"dot"``, ``"dash_dot"``,
        ``"long_dash"``, ``"long_dash_dot"``.
    shadow : dict, optional
        ``{"blur": 3, "offset": 2, "angle": 135,
        "color": RGBColor(0,0,0), "opacity": 0.15}``
        (*blur* / *offset* in points).
    text : str, optional
        Text to render inside the shape.
    font_size : int
        Text size in EMU (default 12 pt).
    color : RGBColor, optional
        Text colour (default ``#404040``).
    bold : bool
        Bold text (default ``False``).
    align : PP_ALIGN
        Horizontal text alignment (default ``CENTER``).
    valign : str
        ``"top"``, ``"middle"``, ``"bottom"`` (default ``"middle"``).

    Returns
    -------
    pptx.shapes.autoshape.Shape
    """
    key = shape_type.lower()
    if key not in _SHAPE_TYPE_MAP:
        raise ValueError(f"Unknown shape_type {shape_type!r}. Supported: {', '.join(sorted(_SHAPE_TYPE_MAP))}")

    shape = slide.shapes.add_shape(
        _SHAPE_TYPE_MAP[key],
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )

    # ── fill ────────────────────────────────────────────────
    if fill is not None:
        fill_rgb = _to_rgb(fill)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if transparency > 0:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                _apply_fill_transparency(spPr, transparency)
    else:
        shape.fill.background()

    # ── line (border) ───────────────────────────────────────
    if line_color is not None:
        line_rgb = _to_rgb(line_color)
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width)
        if line_dash is not None and line_dash in _DASH_MAP:
            shape.line.dash_style = _DASH_MAP[line_dash]
    else:
        shape.line.fill.background()

    # ── shadow ──────────────────────────────────────────────
    if shadow is not None:
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is not None:
            _apply_shadow(spPr, shadow)

    # ── text ────────────────────────────────────────────────
    if text is not None:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", _VALIGN_XML.get(valign, "ctr"))
        p = tf.paragraphs[0]
        p.alignment = align
        text_color = color or _L0_COLOR
        _add_runs(p, text, size=font_size, color=text_color, ascii_bold=bold, unicode_bold=bold)

    return shape


def add_line(
    slide: Slide,
    *,
    left: int,
    top: int,
    width: int,
    height: int = 0,
    color: "Color" = _L0_COLOR,
    line_width: int = 1,
    dash: str | None = None,
):
    """Add a straight line to *slide*.

    Parameters
    ----------
    left, top : int
        Start point in EMU.
    width, height : int
        Horizontal / vertical extent in EMU.
        ``height=0`` → horizontal line, ``width=0`` → vertical line.
    color : RGBColor
        Line colour (default ``#404040``).
    line_width : int
        Thickness in EMU (default 1 pt).
    dash : str, optional
        Dash style (see :func:`add_shape` for options).

    Returns
    -------
    pptx.shapes.connector.Connector
    """
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Emu(left),
        Emu(top),
        Emu(left + width),
        Emu(top + height),
    )
    connector.line.color.rgb = _to_rgb(color)
    connector.line.width = Pt(line_width)
    if dash is not None and dash in _DASH_MAP:
        connector.line.dash_style = _DASH_MAP[dash]
    return connector


# ─── figure ─────────────────────────────────────────────────


def add_figure(
    slide: Slide,
    image_path: str,
    *,
    left: int,
    top: int,
    width: int,
    height: int | None = None,
    caption: str | None = None,
):
    """Add an image to *slide*.

    Parameters
    ----------
    image_path : str
        Path to the image file (PNG, JPEG, etc.).
    left, top : int
        Position of the top-left corner.
    width : int
        Image width.  If *height* is ``None`` the aspect ratio is
        preserved automatically.
    height : int, optional
        Image height.  When omitted, scaled proportionally from *width*.
    caption : str, optional
        Caption text displayed below the image.

    Returns
    -------
    pptx.shapes.picture.Picture
        The inserted picture shape for further customisation.
    """
    pic = slide.shapes.add_picture(
        image_path, Emu(left), Emu(top), Emu(width), Emu(height) if height is not None else None
    )
    if caption:
        _add_caption_below(slide, caption, left=left, bottom=int(pic.top) + int(pic.height), width=width)
    return pic


# ─── table ───────────────────────────────────────────────────

# Default table colours (ground-truth)
_TBL_HEADER_FILL = RGBColor(0xBF, 0xBF, 0xBF)  # 진한 그레이
_TBL_FIRST_COL_FILL = RGBColor(0xE0, 0xE0, 0xE0)  # 옅은 그레이
_TBL_TEXT_COLOR = RGBColor(0x40, 0x40, 0x40)
_TBL_BORDER_COLOR = RGBColor(0x00, 0x00, 0x00)
_TBL_BORDER_WIDTH = 1  # pt


def _set_cell_border(tcPr, color: RGBColor, width: int = _TBL_BORDER_WIDTH) -> None:
    """Add solid borders on all four sides of a table cell."""
    for side in ("lnL", "lnR", "lnT", "lnB"):
        ln = etree.SubElement(tcPr, qn(f"a:{side}"))
        ln.set("w", str(width))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        sf = etree.SubElement(ln, qn("a:solidFill"))
        clr_el = etree.SubElement(sf, qn("a:srgbClr"))
        clr_el.set("val", str(color))
        etree.SubElement(ln, qn("a:prstDash")).set("val", "solid")
        etree.SubElement(ln, qn("a:round"))


def _set_cell_fill(tcPr, color: RGBColor | None) -> None:
    """Set solid fill or noFill on a table cell."""
    if color is None:
        etree.SubElement(tcPr, qn("a:noFill"))
    else:
        sf = etree.SubElement(tcPr, qn("a:solidFill"))
        clr_el = etree.SubElement(sf, qn("a:srgbClr"))
        clr_el.set("val", str(color))


def add_table(
    slide: Slide,
    data: list[list[str]],
    *,
    left: int,
    top: int,
    width: int,
    height: int | None = None,
    col_widths: list[int] | None = None,
    first_row: bool = True,
    first_col: bool = False,
    header_size: int = 14,
    body_size: int = 12,
    header_fill: "Color" = _TBL_HEADER_FILL,
    first_col_fill: "Color" = _TBL_FIRST_COL_FILL,
    text_color: "Color" = _TBL_TEXT_COLOR,
    border_color: "Color" = _TBL_BORDER_COLOR,
    border_width: int = _TBL_BORDER_WIDTH,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    caption: str | None = None,
) -> Table:
    """Add a styled table to *slide*.

    Parameters
    ----------
    data : list[list[str]]
        2-D list of cell text.  ``data[0]`` is the first row
        (header when *first_row* is ``True``).
    left, top, width : int
        Position and width of the table.
    height : int, optional
        Total table height.  Defaults to a sensible auto value.
    col_widths : list, optional
        Per-column widths in EMU.  If ``None``, columns share
        *width* equally.
    first_row : bool
        Fill the first row with *header_fill*.
    first_col : bool
        Fill the first column with *first_col_fill*.
    header_size : int
        Font size for the header row (default 14 pt).
    body_size : int
        Font size for body rows (default 12 pt).
    header_fill : RGBColor
        Background colour of the header row.
    first_col_fill : RGBColor
        Background colour of the first column (body rows only).
    text_color : RGBColor
        Default text colour.
    border_color : RGBColor
        Border colour for all cells.
    border_width : int
        Border thickness in EMU.
    align : PP_ALIGN
        Horizontal text alignment (default ``CENTER``).

    Returns
    -------
    pptx.table.Table
        The created table object for further customisation.
    """
    header_fill_rgb = _to_rgb(header_fill)
    first_col_fill_rgb = _to_rgb(first_col_fill)
    text_color_rgb = _to_rgb(text_color)
    border_color_rgb = _to_rgb(border_color)

    n_rows = len(data)
    n_cols = max(len(row) for row in data) if data else 0
    if n_rows == 0 or n_cols == 0:
        raise ValueError("data must be a non-empty 2-D list")

    row_h = Emu(360000)  # default per-row height
    tbl_height = height or Emu(row_h * n_rows)

    graphic_frame = slide.shapes.add_table(n_rows, n_cols, Emu(left), Emu(top), Emu(width), Emu(tbl_height))
    table = graphic_frame.table

    # ── clear built-in table style to avoid theme conflicts ──
    tbl_el_or_none = graphic_frame._element.find(".//" + qn("a:tbl"))
    assert tbl_el_or_none is not None, "a:tbl element not found"
    tbl_el = tbl_el_or_none
    tblPr = tbl_el.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = etree.SubElement(tbl_el, qn("a:tblPr"))
    # remove tableStyleId — we style cells explicitly
    for child in list(tblPr):
        tblPr.remove(child)
    # disable built-in banding flags
    tblPr.set("firstRow", "0")
    tblPr.set("firstCol", "0")
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")

    # ── column widths ────────────────────────────────────────
    if col_widths:
        for ci, gridCol in enumerate(tbl_el.findall(qn("a:tblGrid") + "/" + qn("a:gridCol"))):
            if ci < len(col_widths):
                gridCol.set("w", str(int(col_widths[ci])))

    # ── populate cells ───────────────────────────────────────
    for ri in range(n_rows):
        is_header = first_row and ri == 0
        font_size = header_size if is_header else body_size

        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell_text = data[ri][ci] if ci < len(data[ri]) else ""

            # ── text ─────────────────────────────────────────
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = align
            _add_runs(p, cell_text, size=font_size, color=text_color_rgb)

            # ── cell properties (fill + borders) ─────────────
            tc = cell._tc
            # remove existing tcPr to rebuild cleanly
            old_tcPr = tc.find(qn("a:tcPr"))
            if old_tcPr is not None:
                tc.remove(old_tcPr)
            tcPr = etree.SubElement(tc, qn("a:tcPr"))
            tcPr.set("anchor", "ctr")

            # borders
            _set_cell_border(tcPr, border_color_rgb, int(Pt(border_width)))

            # fill
            if is_header:
                _set_cell_fill(tcPr, header_fill_rgb)
            elif first_col and ci == 0:
                _set_cell_fill(tcPr, first_col_fill_rgb)
            else:
                _set_cell_fill(tcPr, None)

    if caption:
        _add_caption_below(slide, caption, left=left, bottom=int(top) + int(tbl_height), width=width)

    return table


# ─── chart ───────────────────────────────────────────────────

_CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    # ── category charts ──────────────────────────────────────
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "bar_stacked_100": XL_CHART_TYPE.BAR_STACKED_100,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "line_stacked": XL_CHART_TYPE.LINE_STACKED,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "area_stacked_100": XL_CHART_TYPE.AREA_STACKED_100,
    "radar": XL_CHART_TYPE.RADAR,
    "radar_filled": XL_CHART_TYPE.RADAR_FILLED,
    "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
    # ── pie / doughnut ───────────────────────────────────────
    "pie": XL_CHART_TYPE.PIE,
    "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "doughnut_exploded": XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    # ── scatter (XY) ─────────────────────────────────────────
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
    "scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    "scatter_lines_no_markers": XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    "scatter_smooth_no_markers": XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
    # ── bubble ───────────────────────────────────────────────
    "bubble": XL_CHART_TYPE.BUBBLE,
}

_CATEGORY_TYPES = {
    "bar",
    "bar_stacked",
    "bar_stacked_100",
    "column",
    "column_stacked",
    "column_stacked_100",
    "line",
    "line_markers",
    "line_stacked",
    "area",
    "area_stacked",
    "area_stacked_100",
    "radar",
    "radar_filled",
    "radar_markers",
}
_PIE_TYPES = {"pie", "pie_exploded", "doughnut", "doughnut_exploded"}
_SCATTER_TYPES = {
    "scatter",
    "scatter_lines",
    "scatter_smooth",
    "scatter_lines_no_markers",
    "scatter_smooth_no_markers",
}
_BUBBLE_TYPES = {"bubble"}

# Sensible palette (Office-friendly)
_CHART_PALETTE: list[RGBColor] = [
    RGBColor(0x07, 0x2A, 0x5E),  # navy
    RGBColor(0x00, 0xB0, 0x50),  # green
    RGBColor(0xC0, 0x00, 0x00),  # red
    RGBColor(0xFF, 0x7F, 0x00),  # orange
    RGBColor(0x70, 0x30, 0xA0),  # purple
    RGBColor(0x00, 0x70, 0xC0),  # blue
    RGBColor(0x00, 0x80, 0x80),  # teal
    RGBColor(0x80, 0x80, 0x80),  # gray
]

# Chart text defaults (match template conventions)
_CHART_TITLE_SIZE = 12
_CHART_LABEL_SIZE = 10
_CHART_TEXT_COLOR = RGBColor(0x40, 0x40, 0x40)


def _set_chart_txPr(
    parent,
    *,
    size: int,
    color: RGBColor = _CHART_TEXT_COLOR,
    ascii_font: str = FONT_ASCII,
    ea_font: str = FONT_EA,
    bold: bool = False,
) -> None:
    """Add ``<c:txPr>`` with default run properties (latin + EA fonts).

    Sets the default text style on a chart element (axis, legend, etc.)
    so that PowerPoint renders labels with the template's font conventions.
    """
    for old in parent.findall(qn("c:txPr")):
        parent.remove(old)

    txPr = etree.SubElement(parent, qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))

    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", str(size * 100))  # pt → hundredths-of-pt
    defRPr.set("b", "1" if bold else "0")

    sf = etree.SubElement(defRPr, qn("a:solidFill"))
    srgb = etree.SubElement(sf, qn("a:srgbClr"))
    srgb.set("val", str(color))

    latin = etree.SubElement(defRPr, qn("a:latin"))
    latin.set("typeface", ascii_font)

    ea = etree.SubElement(defRPr, qn("a:ea"))
    ea.set("typeface", ea_font)

    etree.SubElement(p, qn("a:endParaRPr")).set("lang", "ko-KR")


def add_chart(
    slide: Slide,
    chart_type: str,
    *,
    categories: Sequence[str] | None = None,
    series: dict[str, Sequence] | Sequence | None = None,
    left: int,
    top: int,
    width: int,
    height: int,
    title: str | None = None,
    legend: bool = True,
    colors: "Sequence[Color] | None" = None,
    number_format: str | None = None,
    caption: str | None = None,
) -> Chart:
    """Add a chart to *slide*.

    Parameters
    ----------
    chart_type : str
        One of the supported type names (see ``_CHART_TYPE_MAP``).

        **Category** (need *categories* + *series* dict):
        ``"bar"``, ``"bar_stacked"``, ``"bar_stacked_100"``,
        ``"column"``, ``"column_stacked"``, ``"column_stacked_100"``,
        ``"line"``, ``"line_markers"``, ``"line_stacked"``,
        ``"area"``, ``"area_stacked"``, ``"area_stacked_100"``,
        ``"radar"``, ``"radar_filled"``, ``"radar_markers"``

        **Pie / Doughnut** (need *categories* + *series*):
        ``"pie"``, ``"pie_exploded"``,
        ``"doughnut"``, ``"doughnut_exploded"``

        **Scatter** (*series* as ``{name: [(x, y), ...]}``:
        ``"scatter"``, ``"scatter_lines"``, ``"scatter_smooth"``,
        ``"scatter_lines_no_markers"``, ``"scatter_smooth_no_markers"``

        **Bubble** (*series* as ``{name: [(x, y, size), ...]}``:
        ``"bubble"``

    categories : list[str], optional
        Category labels (required for category / pie charts).
    series : dict or list
        Series data.  Format depends on *chart_type*:

        * **Category**: ``{"Series A": [v1, v2, ...], ...}``
        * **Pie**: ``{"Label": [v1, v2, ...]}`` or ``[v1, v2, ...]``
        * **Scatter**: ``{"Series A": [(x1, y1), (x2, y2), ...], ...}``
        * **Bubble**: ``{"Series A": [(x, y, size), ...], ...}``

    title : str, optional
        Chart title text.
    legend : bool
        Show the legend (default ``True``).
    colors : list[RGBColor], optional
        Per-series colours.  Falls back to ``_CHART_PALETTE``.
    number_format : str, optional
        Number format for data labels (e.g. ``"#,##0"``, ``"0.0%"``).

    Returns
    -------
    pptx.chart.chart.Chart
        The created chart for further customisation.
    """
    key = chart_type.lower()
    if key not in _CHART_TYPE_MAP:
        raise ValueError(f"Unknown chart_type {chart_type!r}. Supported: {', '.join(sorted(_CHART_TYPE_MAP))}")

    xl_type = _CHART_TYPE_MAP[key]

    # ── build chart data ─────────────────────────────────────
    if key in _CATEGORY_TYPES or key in _PIE_TYPES:
        if categories is None:
            raise ValueError(f"{chart_type!r} requires 'categories'")
        chart_data = CategoryChartData()
        chart_data.categories = categories
        if isinstance(series, dict):
            for name, values in series.items():
                chart_data.add_series(name, values)
        elif isinstance(series, (list, tuple)):
            # bare list → single unnamed series (pie shorthand)
            chart_data.add_series("", series)
        else:
            raise TypeError("series must be a dict or list for category/pie charts")

    elif key in _SCATTER_TYPES:
        chart_data = XyChartData()
        if not isinstance(series, dict):
            raise TypeError("series must be a dict for scatter charts")
        for name, points in series.items():
            s = chart_data.add_series(name)
            for pt in points:
                s.add_data_point(pt[0], pt[1])

    elif key in _BUBBLE_TYPES:
        chart_data = BubbleChartData()
        if not isinstance(series, dict):
            raise TypeError("series must be a dict for bubble charts")
        for name, points in series.items():
            s = chart_data.add_series(name)
            for pt in points:
                s.add_data_point(pt[0], pt[1], pt[2])
    else:
        raise ValueError(f"Unhandled chart family for {chart_type!r}")

    # ── create chart shape ───────────────────────────────────
    graphic_frame: GraphicFrame = slide.shapes.add_chart(  # type: ignore[assignment]
        xl_type,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
        chart_data,  # type: ignore[arg-type]
    )
    chart = graphic_frame.chart

    # ── title ────────────────────────────────────────────────
    if title is not None:
        chart.has_title = True
        title_tf = chart.chart_title.text_frame
        title_tf.clear()
        _add_runs(title_tf.paragraphs[0], title, size=_CHART_TITLE_SIZE, color=_CHART_TEXT_COLOR)
    else:
        chart.has_title = False

    # ── legend ───────────────────────────────────────────────
    chart.has_legend = legend
    if legend and chart.legend is not None:
        chart.legend.include_in_layout = False
        _set_chart_txPr(chart.legend._element, size=_CHART_LABEL_SIZE)

    # ── axes ─────────────────────────────────────────────────
    if key not in _PIE_TYPES:
        try:
            _set_chart_txPr(chart.category_axis._element, size=_CHART_LABEL_SIZE)
        except (ValueError, AttributeError):
            pass
        try:
            _set_chart_txPr(chart.value_axis._element, size=_CHART_LABEL_SIZE)
        except (ValueError, AttributeError):
            pass

    # ── series / point colours ──────────────────────────────
    palette = [_to_rgb(c) for c in colors] if colors else _CHART_PALETTE
    if key in _PIE_TYPES:
        # Pie/doughnut: colour each slice (data point), not the series
        for s in chart.series:
            for pi, point in enumerate(s.points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = palette[pi % len(palette)]
    else:
        for i, s in enumerate(chart.series):
            if i < len(palette):
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = palette[i % len(palette)]
                # line colour for line/scatter/radar
                if key in _SCATTER_TYPES or key.startswith("line") or key.startswith("radar"):
                    s.format.line.fill.solid()
                    s.format.line.fill.fore_color.rgb = palette[i % len(palette)]

    # ── number format ────────────────────────────────────────
    if number_format is not None:
        for s in chart.series:
            s.number_format = number_format

    # ── caption ──────────────────────────────────────────────
    if caption:
        _add_caption_below(slide, caption, left=left, bottom=int(top) + int(height), width=width)

    return chart


# ─── visual area (auto-layout container) ────────────────────

_VISUAL_TOP_GAP = Emu(100_000)  # gap between content box bottom and visual area top
_VISUAL_GAP = Emu(150_000)  # horizontal gap between elements (~0.17")
_VISUAL_BOTTOM_MARGIN = Emu(258_000)  # bottom margin from slide edge


class VisualArea:
    """Auto-layout container that arranges visual elements in a horizontal row.

    Elements are collected via :meth:`add_chart`, :meth:`add_table`,
    :meth:`add_figure`, and :meth:`add_shape`, then placed all at once
    when :meth:`render` is called.
    """

    def __init__(
        self,
        slide: Slide,
        *,
        left: int,
        top: int,
        width: int,
        height: int,
        gap: int,
    ) -> None:
        self._slide = slide
        self._left = int(left)
        self._top = int(top)
        self._width = int(width)
        self._height = int(height)
        self._gap = int(gap)
        self._elements: list[dict] = []
        self._rendered = False

    # ── element collectors ──────────────────────────────────

    def _guard_add(self, weight: float) -> None:
        if self._rendered:
            raise RuntimeError("render() 호출 후에는 요소를 추가할 수 없습니다.")
        if weight <= 0:
            raise ValueError("weight는 0보다 커야 합니다.")

    def add_chart(
        self,
        chart_type: str,
        *,
        weight: float = 1.0,
        categories: Sequence[str] | None = None,
        series: dict[str, Sequence] | Sequence | None = None,
        title: str | None = None,
        legend: bool = True,
        colors: "Sequence[Color] | None" = None,
        number_format: str | None = None,
        caption: str | None = None,
    ) -> "VisualArea":
        """Add a chart element to the container."""
        self._guard_add(weight)
        self._elements.append({
            "kind": "chart",
            "weight": weight,
            "kwargs": {
                "chart_type": chart_type,
                "categories": categories,
                "series": series,
                "title": title,
                "legend": legend,
                "colors": colors,
                "number_format": number_format,
                "caption": caption,
            },
        })
        return self

    def add_table(
        self,
        data: list[list[str]],
        *,
        weight: float = 1.0,
        col_widths: list[int] | None = None,
        first_row: bool = True,
        first_col: bool = False,
        header_size: int = 14,
        body_size: int = 12,
        header_fill: "Color" = _TBL_HEADER_FILL,
        first_col_fill: "Color" = _TBL_FIRST_COL_FILL,
        text_color: "Color" = _TBL_TEXT_COLOR,
        border_color: "Color" = _TBL_BORDER_COLOR,
        border_width: int = _TBL_BORDER_WIDTH,
        align: PP_ALIGN = PP_ALIGN.CENTER,
        caption: str | None = None,
    ) -> "VisualArea":
        """Add a table element to the container."""
        self._guard_add(weight)
        self._elements.append({
            "kind": "table",
            "weight": weight,
            "kwargs": {
                "data": data,
                "col_widths": col_widths,
                "first_row": first_row,
                "first_col": first_col,
                "header_size": header_size,
                "body_size": body_size,
                "header_fill": header_fill,
                "first_col_fill": first_col_fill,
                "text_color": text_color,
                "border_color": border_color,
                "border_width": border_width,
                "align": align,
                "caption": caption,
            },
        })
        return self

    def add_figure(
        self,
        image_path: str,
        *,
        weight: float = 1.0,
        caption: str | None = None,
    ) -> "VisualArea":
        """Add an image element to the container."""
        self._guard_add(weight)
        self._elements.append({
            "kind": "figure",
            "weight": weight,
            "kwargs": {
                "image_path": image_path,
                "caption": caption,
            },
        })
        return self

    def add_shape(
        self,
        shape_type: str,
        *,
        weight: float = 1.0,
        fill: "Color | None" = None,
        transparency: int = 0,
        line_color: "Color | None" = None,
        line_width: int = 1,
        line_dash: str | None = None,
        shadow: dict | None = None,
        text: str | None = None,
        font_size: int = 12,
        color: "Color | None" = None,
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.CENTER,
        valign: str = "middle",
    ) -> "VisualArea":
        """Add a shape element to the container."""
        self._guard_add(weight)
        self._elements.append({
            "kind": "shape",
            "weight": weight,
            "kwargs": {
                "shape_type": shape_type,
                "fill": fill,
                "transparency": transparency,
                "line_color": line_color,
                "line_width": line_width,
                "line_dash": line_dash,
                "shadow": shadow,
                "text": text,
                "font_size": font_size,
                "color": color,
                "bold": bold,
                "align": align,
                "valign": valign,
            },
        })
        return self

    # ── rendering ───────────────────────────────────────────

    def render(self) -> list:
        """Calculate positions and place all elements on the slide.

        Returns a list of placed objects (Chart, Table, Picture, Shape).
        Can only be called once.
        """
        if self._rendered:
            raise RuntimeError("VisualArea.render()는 한 번만 호출할 수 있습니다.")
        self._rendered = True

        n = len(self._elements)
        if n == 0:
            return []

        caption_reserve = int(_CAPTION_GAP) + int(_CAPTION_BOX_HEIGHT)

        total_weight = sum(e["weight"] for e in self._elements)
        total_gap = self._gap * (n - 1)
        available_width = self._width - total_gap

        results: list = []
        cursor_left = self._left

        for i, elem in enumerate(self._elements):
            if i == n - 1:
                elem_width = (self._left + self._width) - cursor_left
            else:
                elem_width = int(available_width * (elem["weight"] / total_weight))

            if elem["kwargs"].get("caption"):
                elem_height = self._height - caption_reserve
            else:
                elem_height = self._height

            result = self._place_element(elem, cursor_left, self._top, elem_width, elem_height)
            results.append(result)
            cursor_left += elem_width + self._gap

        return results

    def _place_element(self, elem: dict, left: int, top: int, width: int, height: int):
        """Dispatch to the appropriate ``add_*`` function."""
        kind = elem["kind"]
        kw = elem["kwargs"]

        if kind == "chart":
            return add_chart(
                self._slide,
                kw["chart_type"],
                categories=kw["categories"],
                series=kw["series"],
                left=left,
                top=top,
                width=width,
                height=height,
                title=kw["title"],
                legend=kw["legend"],
                colors=kw["colors"],
                number_format=kw["number_format"],
                caption=kw["caption"],
            )
        elif kind == "table":
            return add_table(
                self._slide,
                kw["data"],
                left=left,
                top=top,
                width=width,
                col_widths=kw["col_widths"],
                first_row=kw["first_row"],
                first_col=kw["first_col"],
                header_size=kw["header_size"],
                body_size=kw["body_size"],
                header_fill=kw["header_fill"],
                first_col_fill=kw["first_col_fill"],
                text_color=kw["text_color"],
                border_color=kw["border_color"],
                border_width=kw["border_width"],
                align=kw["align"],
                caption=kw["caption"],
            )
        elif kind == "figure":
            return add_figure(
                self._slide,
                kw["image_path"],
                left=left,
                top=top,
                width=width,
                caption=kw["caption"],
            )
        elif kind == "shape":
            return add_shape(
                self._slide,
                kw["shape_type"],
                left=left,
                top=top,
                width=width,
                height=height,
                fill=kw["fill"],
                transparency=kw["transparency"],
                line_color=kw["line_color"],
                line_width=kw["line_width"],
                line_dash=kw["line_dash"],
                shadow=kw["shadow"],
                text=kw["text"],
                font_size=kw["font_size"],
                color=kw["color"],
                bold=kw["bold"],
                align=kw["align"],
                valign=kw["valign"],
            )
        else:
            raise ValueError(f"Unknown element kind: {kind!r}")


def visual_area(
    slide: Slide,
    *,
    left: int = _CONTENT_LEFT,
    top: int | None = None,
    width: int = _CONTENT_WIDTH,
    height: int | None = None,
    gap: int = _VISUAL_GAP,
) -> VisualArea:
    """Create a :class:`VisualArea` for auto-layout of visual elements.

    Parameters
    ----------
    slide : Slide
        Target slide.
    left : int
        Left edge in EMU (default: content box left).
    top : int, optional
        Top edge in EMU.  Defaults to content box bottom + gap.
    width : int
        Width in EMU (default: content box width).
    height : int, optional
        Height in EMU.  Defaults to fill to bottom margin.
    gap : int
        Horizontal gap between elements (default ~0.17").
    """
    if top is None:
        top = int(_CONTENT_TOP) + int(_CONTENT_HEIGHT) + int(_VISUAL_TOP_GAP)
    if height is None:
        height = 6_858_000 - int(_VISUAL_BOTTOM_MARGIN) - int(top)
    return VisualArea(
        slide, left=int(left), top=int(top), width=int(width), height=int(height), gap=int(gap)
    )


# ─── cover slide helpers ─────────────────────────────────────


def _set_cover_pPr(pPr, *, spc_before: int = 0) -> None:
    """Apply ground-truth paragraph properties for cover shapes."""
    pPr.set("marL", "0")
    pPr.set("marR", "0")
    pPr.set("indent", "0")
    pPr.set("eaLnBrk", "1")
    pPr.set("fontAlgn", "base")
    pPr.set("latinLnBrk", "1")
    pPr.set("hangingPunct", "1")
    # 150 % line spacing
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct")).set("val", "150000")
    # space before
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPct")).set("val", str(spc_before))
    # space after = 0
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPct")).set("val", "0")
    # no bullet
    etree.SubElement(pPr, qn("a:buNone"))


# ─── cover slide API ────────────────────────────────────────


def add_cover_title(
    slide: Slide,
    text: str,
    *,
    left: int = _COVER_TITLE_LEFT,
    top: int = _COVER_TITLE_TOP,
    width: int = _COVER_TITLE_WIDTH,
    height: int = _COVER_TITLE_HEIGHT,
    font_size: int = _COVER_TITLE_SIZE,
    color: "Color" = _COVER_COLOR,
) -> None:
    """Add the centred presentation title (cover slides only).

    Use ``\\n`` in *text* for line breaks.  Any number of consecutive
    newlines is normalised to exactly **two** soft returns so the
    visual gap is always consistent.
    """
    _require(slide, _COVER)
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    # match ground-truth: auto-fit + word-wrap
    bodyPr = _find_elem(tf._txBody, "a:bodyPr")
    bodyPr.set("wrap", "square")
    for old in bodyPr.findall(qn("a:spAutoFit")) + bodyPr.findall(qn("a:noAutofit")):
        bodyPr.remove(old)
    etree.SubElement(bodyPr, qn("a:spAutoFit"))
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    pPr = p._element.get_or_add_pPr()
    _set_cover_pPr(pPr, spc_before=50000)
    parts = [s for s in re.split(r"\n+", text) if s]
    for i, part in enumerate(parts):
        if i > 0:
            sz_full = str(font_size * 100)  # pt → hundredths-of-pt
            sz_half = str(font_size * 50)  # half-size for tighter gap
            for sz_val in (sz_full, sz_half):
                br = etree.SubElement(p._element, qn("a:br"))
                brPr = etree.SubElement(br, qn("a:rPr"))
                brPr.set("lang", "en-US")
                brPr.set("sz", sz_val)
        _add_runs(p, part, size=font_size, color=color)


def add_cover_info(
    slide: Slide,
    date: str,
    presenter: str,
    *,
    left: int = _COVER_INFO_LEFT,
    top: int = _COVER_INFO_TOP,
    width: int = _COVER_INFO_WIDTH,
    height: int = _COVER_INFO_HEIGHT,
    font_size: int = _COVER_INFO_SIZE,
    color: "Color" = _COVER_COLOR,
) -> None:
    """Add date and presenter info (cover slides only)."""
    _require(slide, _COVER)
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    bodyPr = _find_elem(tf._txBody, "a:bodyPr")
    bodyPr.set("wrap", "square")
    for old in bodyPr.findall(qn("a:spAutoFit")) + bodyPr.findall(qn("a:noAutofit")):
        bodyPr.remove(old)
    etree.SubElement(bodyPr, qn("a:spAutoFit"))
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_cover_pPr(p._element.get_or_add_pPr())
    _add_runs(p, f"날짜 : {date}", size=font_size, color=color)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    _set_cover_pPr(p._element.get_or_add_pPr())
    _add_runs(p, f"발표자 : {presenter}", size=font_size, color=color)


def _find_elem(parent: etree._Element, tag: str) -> etree._Element:
    """Find the first child element with the given tag, or return None."""
    elem = parent.find(qn(tag))
    if elem is None:
        raise ValueError(f"Expected element '{tag}' not found")
    return elem
