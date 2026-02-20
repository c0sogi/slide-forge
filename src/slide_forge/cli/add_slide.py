"""Add or duplicate slides in an existing PPTX file.

Usage:
    slide-forge add-slide presentation.pptx
    slide-forge add-slide presentation.pptx --type cover
    slide-forge add-slide presentation.pptx --source 3
    slide-forge add-slide presentation.pptx --source 3 --output updated.pptx
"""

from __future__ import annotations

import argparse
import sys
from copy import deepcopy
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation

from slide_forge import create_cover_slide, create_slide

if TYPE_CHECKING:
    from pptx.presentation import Presentation as _Presentation

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def configure_parser(subparsers: argparse._SubParsersAction) -> None:
    parser = subparsers.add_parser("add-slide", help="Add or duplicate a slide in a PPTX file")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument(
        "--type",
        choices=["content", "cover"],
        default="content",
        dest="slide_type",
        help="Slide type to add (default: content). Ignored when --source is used.",
    )
    parser.add_argument(
        "--source",
        type=int,
        default=None,
        metavar="N",
        help="Duplicate slide N (1-based) instead of adding a blank slide",
    )
    parser.add_argument(
        "--output",
        default=None,
        help="Output file path (default: overwrite input file)",
    )
    parser.set_defaults(func=_run)


def _duplicate_slide(prs: _Presentation, source_index: int) -> None:
    """Duplicate the slide at 1-based *source_index*, appending the copy at the end."""
    slides = prs.slides
    if source_index < 1 or source_index > len(slides):
        print(
            f"Error: slide {source_index} out of range (1-{len(slides)})",
            file=sys.stderr,
        )
        sys.exit(1)

    source = slides[source_index - 1]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # ── Replace new slide XML with deep copy of source ──
    new_el = new_slide._element
    for child in list(new_el):
        new_el.remove(child)
    for child in source._element:
        new_el.append(deepcopy(child))

    # ── Copy relationships and build rId mapping ──
    rId_map: dict[str, str] = {}
    src_rels = source.part._rels
    new_rels = new_slide.part._rels

    # Map the layout rId (already exists on new slide)
    for src_rId in src_rels:
        if "slideLayout" in src_rels[src_rId].reltype:
            for new_rId in new_rels:
                if "slideLayout" in new_rels[new_rId].reltype:
                    rId_map[src_rId] = new_rId
                    break
            break

    # Copy other relationships (images, charts, hyperlinks, etc.)
    for src_rId in src_rels:
        src_rel = src_rels[src_rId]
        if "slideLayout" in src_rel.reltype or "notesSlide" in src_rel.reltype:
            continue
        if src_rel.is_external:
            new_rId = new_slide.part.relate_to(
                src_rel.target_ref, src_rel.reltype, is_external=True
            )
        else:
            new_rId = new_slide.part.relate_to(src_rel.target_part, src_rel.reltype)
        rId_map[src_rId] = new_rId

    # ── Update rId references in copied XML ──
    r_attrs = (f"{{{_R_NS}}}embed", f"{{{_R_NS}}}link", f"{{{_R_NS}}}id")
    for elem in new_el.iter():
        for attr in r_attrs:
            val = elem.get(attr)
            if val and val in rId_map and rId_map[val] != val:
                elem.set(attr, rId_map[val])


def _run(args: argparse.Namespace) -> None:
    pptx_path = Path(args.pptx)

    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    if pptx_path.suffix.lower() != ".pptx":
        print(f"Error: {pptx_path} must be a .pptx file", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))

    if args.source is not None:
        _duplicate_slide(prs, args.source)
        print(f"Duplicated slide {args.source} in {pptx_path.name}")
    elif args.slide_type == "cover":
        create_cover_slide(prs)
        print(f"Added cover slide to {pptx_path.name}")
    else:
        create_slide(prs)
        print(f"Added content slide to {pptx_path.name}")

    output_path = Path(args.output) if args.output else pptx_path
    prs.save(str(output_path))
    print(f"Saved to {output_path}")
